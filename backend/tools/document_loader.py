"""
JARVIS — Document Loader
Unified file parser supporting PDF, TXT, DOCX, PPTX, and images.
Extracted from analyse.py for reuse across agents.
"""

import base64
import os

from langchain_community.document_loaders import PyPDFLoader
from backend.logger import get_logger

logger = get_logger("tools.document_loader")


def encode_image(image_path: str) -> str:
    """Encode an image file to base64 string."""
    with open(image_path, "rb") as image_file:
        return base64.b64encode(image_file.read()).decode("utf-8")


def load_and_parse_file(file_path: str, vision_llm=None) -> str:
    """
    Parse a file and return its text content.

    Supports: .txt, .md, .pdf, .docx, .pptx, .png, .jpg, .jpeg
    For images, requires a vision_llm to be passed.
    """
    ext = file_path.lower().rsplit(".", 1)[-1] if "." in file_path else ""
    logger.info(f"Parsing file: {os.path.basename(file_path)} (type: .{ext})")

    if ext in ("txt", "md"):
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()

    elif ext == "pdf":
        loader = PyPDFLoader(file_path)
        docs = loader.load()
        return "\n".join([doc.page_content for doc in docs])

    elif ext == "docx":
        import docx2txt
        from docx import Document

        text = docx2txt.process(file_path)
        try:
            doc = Document(file_path)
            tables_text = []
            for table in doc.tables:
                table_data = []
                for row in table.rows:
                    row_data = [cell.text.strip() for cell in row.cells]
                    table_data.append(" | ".join(row_data))
                tables_text.append("\n".join(table_data))
            if tables_text:
                text += "\n\nExtracted Tables:\n" + "\n\n".join(tables_text)
        except Exception as e:
            logger.warning(f"Error extracting Word tables: {e}")
        return text

    elif ext == "pptx":
        from pptx import Presentation

        prs = Presentation(file_path)
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_runs.append(shape.text.strip())
                if shape.has_table:
                    table_data = []
                    for row in shape.table.rows:
                        row_data = [cell.text.strip() for cell in row.cells]
                        table_data.append(" | ".join(row_data))
                    text_runs.append("Slide Table:\n" + "\n".join(table_data))
        return "\n".join(text_runs)

    elif ext in ("png", "jpg", "jpeg"):
        if not vision_llm:
            logger.warning("No vision LLM provided for image parsing.")
            return "[Image file — vision LLM required for extraction]"

        base64_image = encode_image(file_path)
        prompt = [
            {
                "role": "user",
                "content": [
                    {
                        "type": "text",
                        "text": "Extract all text, tables, and describe the content of this image in detail for indexing in a database.",
                    },
                    {
                        "type": "image_url",
                        "image_url": {"url": f"data:image/jpeg;base64,{base64_image}"},
                    },
                ],
            }
        ]
        response = vision_llm.invoke(prompt)
        return response.content

    elif ext in ("mp4", "mkv", "avi", "mov", "webm"):
        if not vision_llm:
            logger.warning("No vision LLM provided for video parsing.")
            return "[Video file — vision LLM required for extraction]"

        import cv2
        cap = cv2.VideoCapture(file_path)
        if not cap.isOpened():
            logger.error(f"Failed to open video: {file_path}")
            return "[Error opening video file]"

        fps = cap.get(cv2.CAP_PROP_FPS)
        total_frames = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))
        duration = total_frames / fps if fps > 0 else 0
        logger.info(f"Video details: {fps} FPS, {total_frames} total frames, duration: {duration:.2f}s")

        num_frames_to_extract = min(10, total_frames) if total_frames > 0 else 0
        if num_frames_to_extract == 0:
            cap.release()
            return "[Video contains no frames]"

        indices = [int(i * total_frames / num_frames_to_extract) for i in range(num_frames_to_extract)]
        frame_descriptions = []

        for idx, frame_idx in enumerate(indices):
            cap.set(cv2.CAP_PROP_POS_FRAMES, frame_idx)
            ret, frame = cap.read()
            if not ret:
                continue
            
            timestamp_sec = frame_idx / fps if fps > 0 else 0
            timestamp_str = f"{int(timestamp_sec // 60):02d}:{int(timestamp_sec % 60):02d}"

            h, w = frame.shape[:2]
            max_size = 640
            if max(h, w) > max_size:
                scale = max_size / max(h, w)
                frame = cv2.resize(frame, (int(w * scale), int(h * scale)))

            success, buffer = cv2.imencode('.jpg', frame)
            if not success:
                continue
            
            base64_image = base64.b64encode(buffer).decode("utf-8")
            
            prompt = [
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "text",
                            "text": f"This is a frame at timestamp {timestamp_str} from an uploaded video. Briefly describe what is visible in this frame, focusing on actions, text, objects, and setting.",
                        },
                        {
                            "type": "image_url",
                            "image_url": {"url": f"data:image/jpeg;base64,{base64_image}"},
                        },
                    ],
                }
            ]
            try:
                response = vision_llm.invoke(prompt)
                frame_desc = response.content.strip()
                frame_descriptions.append(f"At {timestamp_str}: {frame_desc}")
                logger.info(f"Analyzed frame at {timestamp_str}")
            except Exception as e:
                logger.error(f"Error analyzing frame at {timestamp_str}: {e}")

        cap.release()

        if not frame_descriptions:
            return "[Video parsing failed to extract/analyze any frames]"

        from backend.config import llm
        
        full_timeline = "\n".join(frame_descriptions)
        summary_prompt = (
            f"You are analyzing a video's visual timeline. Below are the descriptions of key frames extracted from the video:\n\n"
            f"{full_timeline}\n\n"
            f"Please synthesize these frame descriptions into a detailed, coherent summary of the video. "
            f"Include the general topic/theme, the progression of events/actions over time, and any key text or objects observed. "
            f"Keep it structured and descriptive for database search indexing."
        )
        
        try:
            summary_response = llm.invoke(summary_prompt)
            video_summary = summary_response.content.strip()
            return f"Video Title/Source: {os.path.basename(file_path)}\n\n--- VIDEO SUMMARY ---\n{video_summary}\n\n--- DETAILED TIMELINE ---\n{full_timeline}"
        except Exception as e:
            logger.error(f"Failed to generate summary: {e}")
            return f"Video Title/Source: {os.path.basename(file_path)}\n\n--- DETAILED TIMELINE ---\n{full_timeline}"

    elif ext == "csv":
        import pandas as pd
        logger.info(f"Parsing CSV file with pandas: {file_path}")
        try:
            df = pd.read_csv(file_path)
            max_rows = 150
            truncated = len(df) > max_rows
            df_subset = df.head(max_rows)
            try:
                md_table = df_subset.to_markdown(index=False)
            except Exception:
                md_table = df_subset.to_string(index=False)
            summary = f"CSV File Structure and Sample (showing first {max_rows} rows):\n"
            summary += f"Columns: {', '.join(map(str, df.columns))}\n"
            summary += f"Total rows: {len(df)}, Total columns: {len(df.columns)}\n\n"
            summary += md_table
            if truncated:
                summary += "\n\n*(Note: Dataset was truncated for LLM analysis)*"
            return summary
        except Exception as e:
            logger.error(f"Error parsing CSV: {e}")
            return f"[Error parsing CSV: {str(e)}]"

    elif ext in ("xlsx", "xls"):
        import pandas as pd
        logger.info(f"Parsing Excel file with pandas: {file_path}")
        try:
            xls = pd.ExcelFile(file_path)
            sheets_summary = []
            for sheet_name in xls.sheet_names:
                df = pd.read_excel(file_path, sheet_name=sheet_name)
                max_rows = 50
                truncated = len(df) > max_rows
                df_subset = df.head(max_rows)
                try:
                    md_table = df_subset.to_markdown(index=False)
                except Exception:
                    md_table = df_subset.to_string(index=False)
                sheet_sum = f"Sheet Name: {sheet_name}\n"
                sheet_sum += f"Columns: {', '.join(map(str, df.columns))}\n"
                sheet_sum += f"Total rows: {len(df)}, Total columns: {len(df.columns)}\n\n"
                sheet_sum += md_table
                if truncated:
                    sheet_sum += "\n\n*(Note: Sheet was truncated for LLM analysis)*"
                sheets_summary.append(sheet_sum)
            return "\n\n===\n\n".join(sheets_summary)
        except Exception as e:
            logger.error(f"Error parsing Excel: {e}")
            return f"[Error parsing Excel: {str(e)}]"

    else:
        logger.warning(f"Unsupported file type: .{ext}")
        return ""
